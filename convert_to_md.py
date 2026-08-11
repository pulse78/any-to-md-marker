#!/usr/bin/env python3
"""
convert_to_md.py — any-to-md-marker (Fork mit marker-pdf-Engine)
===============================================================

Hochpräzise Konvertierung von PDFs (und anderen Formaten) zu Markdown.

Engines für PDF (--engine):
  marker   marker-pdf (Standard) — OCR, Tabellen, Reading Order, Formeln,
           Code-Blöcke, Header/Footer-Entfernung. Höchste Präzision.
           Benötigt ~3-5 GB (PyTorch + Modelle, Download beim ersten Lauf).
  classic  Bisherige Logik: pdfplumber → pymupdf → Tesseract OCR.
  auto     marker, wenn installiert; sonst classic.

Alle anderen Formate (CSV, DOCX, PPTX, HTML, Bilder, Audio) sind unverändert.

Pip-Abhängigkeiten werden beim ersten Bedarf automatisch installiert.

Systemabhängigkeiten (einmalig manuell installieren):
  Tesseract OCR (für Bild-PDFs und Bilddateien):
    macOS:  brew install tesseract tesseract-lang
    Ubuntu: sudo apt install tesseract-ocr tesseract-ocr-deu tesseract-ocr-eng
    Windows: https://github.com/UB-Mannheim/tesseract/wiki

  ffmpeg (für Audio-Transkription):
    macOS:  brew install ffmpeg
    Ubuntu: sudo apt install ffmpeg
    Windows: https://ffmpeg.org/download.html

Usage:
    python convert_to_md.py <eingabe> [ausgabe.md]
        [--engine marker|classic|auto]
        [--encoding <enc>]
        [--lang <tesseract-lang>]   z.B. deu+eng
        [--model <whisper-modell>]  tiny | base | small | medium | large (Standard: base)

Batch (ganzer Ordner):
    python convert_to_md.py <ordner> --batch <ausgabe-ordner>
        [--engine ...]   Konvertiert alle PDFs/unterstützten Dateien im Ordner.

Wenn ausgabe.md fehlt, wird das Ergebnis auf stdout ausgegeben.
"""

import sys
import subprocess
import shutil
from pathlib import Path


# ─────────────────────────────────────────────────────────────────────────────
# Abhängigkeitsverwaltung
# ─────────────────────────────────────────────────────────────────────────────

def pip_install(package: str):
    """Installiert ein pip-Paket automatisch."""
    subprocess.check_call(
        [sys.executable, "-m", "pip", "install", package,
         "--quiet", "--break-system-packages"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def ensure(package: str, import_name: str | None = None):
    """Stellt sicher, dass ein pip-Paket importierbar ist; installiert es sonst."""
    name = import_name or package.split("[")[0].replace("-", "_")
    try:
        __import__(name)
    except ImportError:
        print(f"  Installiere {package} …", file=sys.stderr)
        pip_install(package)


def ensure_ffmpeg() -> str:
    """Prüft ob ffmpeg verfügbar ist (wird von Whisper benötigt)."""
    path = shutil.which("ffmpeg")
    if path:
        return path
    msg = (
        "\n  ✗ ffmpeg nicht gefunden. Bitte installieren:\n"
        "    macOS  : brew install ffmpeg\n"
        "    Ubuntu : sudo apt install ffmpeg\n"
        "    Windows: https://ffmpeg.org/download.html\n"
    )
    print(msg, file=sys.stderr)
    sys.exit(1)


def ensure_tesseract() -> str:
    """
    Prüft ob das Tesseract-Binary verfügbar ist.
    Gibt den Pfad zurück oder bricht mit einer Installationsanleitung ab.
    """
    path = shutil.which("tesseract")
    if path:
        return path
    msg = (
        "\n  ✗ Tesseract nicht gefunden. Bitte installieren:\n"
        "    macOS  : brew install tesseract tesseract-lang\n"
        "    Ubuntu : sudo apt install tesseract-ocr tesseract-ocr-deu tesseract-ocr-eng\n"
        "    Windows: https://github.com/UB-Mannheim/tesseract/wiki\n"
    )
    print(msg, file=sys.stderr)
    sys.exit(1)


def best_ocr_lang() -> str:
    """Wählt die beste verfügbare Tesseract-Sprachkombination."""
    try:
        result = subprocess.run(
            ["tesseract", "--list-langs"],
            capture_output=True, text=True
        )
        available = set(result.stdout.strip().splitlines()[1:])
        if "deu" in available and "eng" in available:
            return "deu+eng"
        elif "deu" in available:
            return "deu"
        return "eng"
    except Exception:
        return "eng"


def marker_available() -> bool:
    """Prüft, ob marker-pdf erreichbar ist (isolierte venv neben dem Skript)."""
    return _marker_single_path() is not None


def _marker_single_path():
    """
    Liefert den Pfad zur marker_single-Executable der isolierten venv.
    Sucht neben dem Skript: <script_dir>/venv/Scripts/marker_single(.exe)
    """
    script_dir = Path(__file__).resolve().parent
    candidates = [
        script_dir / "venv" / "Scripts" / "marker_single.exe",
        script_dir / "venv" / "Scripts" / "marker_single",
        script_dir / "venv" / "bin" / "marker_single",
    ]
    for c in candidates:
        if c.exists():
            return c
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Format-Konverter
# ─────────────────────────────────────────────────────────────────────────────

def via_csv(p: Path) -> str:
    ensure("pandas")
    ensure("tabulate")        # wird von pandas.to_markdown() benötigt
    import pandas as pd
    df = pd.read_csv(p, encoding="utf-8-sig")
    return df.to_markdown(index=False)


def via_xlsx(p: Path) -> str:
    ensure("openpyxl")
    ensure("pandas")
    ensure("tabulate")
    import pandas as pd
    xl = pd.ExcelFile(p, engine="openpyxl")
    sections = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        sections.append(f"## Sheet: {sheet}\n\n{df.to_markdown(index=False)}")
    return "\n\n---\n\n".join(sections)


def via_html(p: Path) -> str:
    ensure("beautifulsoup4", "bs4")
    ensure("markdownify")
    from bs4 import BeautifulSoup
    import markdownify
    html = p.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return markdownify.markdownify(str(soup), heading_style="ATX")


def via_image_ocr(p: Path, lang: str) -> str:
    """
    OCR für Bilddateien (PNG, JPG, TIFF, …) via Tesseract.
    Gibt den erkannten Text als einfachen Markdown-Block zurück.
    """
    ensure("Pillow", "PIL")
    ensure("pytesseract")
    ensure_tesseract()

    from PIL import Image
    import pytesseract

    img = Image.open(p)
    text = pytesseract.image_to_string(img, lang=lang).strip()
    if not text:
        return "*[Kein Text erkannt – Bild möglicherweise zu niedrig aufgelöst oder leer.]*"
    return text


def via_pdf_marker(p: Path) -> str:
    """
    Hochpräzise PDF-Konvertierung via marker-pdf (isolierte venv).
    Nutzt die marker_single-Executable der venv; Ausgabe-Markdown wird
    eingelesen und zurückgegeben.
    """
    marker_single = _marker_single_path()
    if marker_single is None:
        raise RuntimeError("marker_single nicht gefunden (venv fehlt?)")

    import tempfile
    import os
    with tempfile.TemporaryDirectory(prefix="mdmarker_") as tmp:
        out_dir = Path(tmp)
        cmd = [str(marker_single), str(p), "--output_dir", str(out_dir)]
        # PYTHONPATH bereinigen: verhindert, dass marker_single Pakete
        # aus der Haupt-venv (z.B. Hermes) lädt.
        env = os.environ.copy()
        env.pop("PYTHONPATH", None)
        print(f"  marker_single: {p.name} …", file=sys.stderr)
        result = subprocess.run(cmd, capture_output=True, text=True, env=env)
        if result.returncode != 0:
            raise RuntimeError(
                f"marker_single fehlgeschlagen (RC {result.returncode}): "
                f"{result.stderr[-500:]}"
            )
        md_files = sorted(out_dir.rglob("*.md"))
        if not md_files:
            raise RuntimeError("marker_single erzeugte keine .md-Ausgabe")
        return md_files[0].read_text(encoding="utf-8", errors="replace")


def via_pdf(p: Path, lang: str) -> str:
    """
    PDF-Konvertierung (klassische Pipeline):
      1. Versucht Text-Layer mit pdfplumber zu lesen.
      2. Fallback: PyMuPDF.
      3. Wenn kein Text gefunden (Bild-PDF / Scan): OCR via Tesseract.
    """
    pages = []

    # — Versuch 1: pdfplumber ------------------------------------------------
    try:
        ensure("pdfplumber")
        import pdfplumber
        with pdfplumber.open(p) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = (page.extract_text() or "").strip()
                if text:
                    pages.append((i, text))
    except Exception:
        pass

    # — Versuch 2: PyMuPDF ---------------------------------------------------
    if not pages:
        try:
            ensure("pymupdf", "fitz")
            import fitz
            doc = fitz.open(str(p))
            for i, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    pages.append((i, text))
        except Exception:
            pass

    # — Versuch 3: Tesseract OCR (Bild-PDF / Scan) ---------------------------
    if not pages:
        print(
            "  Kein Textlayer gefunden – starte Tesseract OCR …",
            file=sys.stderr
        )
        ensure("pdf2image")
        ensure("pytesseract")
        ensure("Pillow", "PIL")
        ensure_tesseract()

        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(str(p), dpi=300)
        for i, img in enumerate(images, 1):
            text = pytesseract.image_to_string(img, lang=lang).strip()
            if text:
                pages.append((i, text))

        if not pages:
            return "*[Kein Text erkannt – PDF enthält möglicherweise nur nicht-textuellen Inhalt.]*"

    return "\n\n---\n\n".join(f"## Seite {i}\n\n{t}" for i, t in pages)


def via_docx(p: Path) -> str:
    ensure("python-docx", "docx")
    from docx import Document
    doc = Document(str(p))
    lines = []
    numbered_counters: dict[int, int] = {}

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text  = para.text.strip()
        if not text:
            lines.append("")
            continue

        if style == "Title":
            lines.append(f"# {text}")
        elif "Heading 1" in style:
            lines.append(f"## {text}")
        elif "Heading 2" in style:
            lines.append(f"### {text}")
        elif "Heading 3" in style:
            lines.append(f"#### {text}")
        elif "List Bullet" in style:
            try:
                lvl = para._p.pPr.numPr.ilvl.val if (
                    para._p.pPr is not None and
                    para._p.pPr.numPr is not None
                ) else 0
            except Exception:
                lvl = 0
            lines.append(f"{'  ' * lvl}- {text}")
        elif "List Number" in style:
            try:
                lvl = para._p.pPr.numPr.ilvl.val if (
                    para._p.pPr is not None and
                    para._p.pPr.numPr is not None
                ) else 0
            except Exception:
                lvl = 0
            numbered_counters[lvl] = numbered_counters.get(lvl, 0) + 1
            lines.append(f"{'  ' * lvl}{numbered_counters[lvl]}. {text}")
        else:
            lines.append(text)

    return "\n".join(lines)


def via_pptx(p: Path) -> str:
    ensure("python-pptx", "pptx")
    from pptx import Presentation
    prs = Presentation(str(p))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"## Folie {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides.append("\n\n".join(parts))
    return "\n\n---\n\n".join(slides)


def via_plain(p: Path, encoding: str = "utf-8") -> str:
    return p.read_text(encoding=encoding, errors="replace")


def via_audio(p: Path, model_name: str = "base") -> str:
    """
    Transkription von Audio-/Videodateien via OpenAI Whisper (lokal, kein API-Key).

    Beim ersten Aufruf wird das Whisper-Modell heruntergeladen:
      tiny   ~  75 MB  – schnell, niedrigste Qualität
      base   ~ 150 MB  – guter Kompromiss (Standard)
      small  ~ 500 MB  – bessere Qualität
      medium ~ 1,5 GB  – sehr gut
      large  ~ 3,0 GB  – beste Qualität, langsam

    Das Modell wird im Home-Verzeichnis gecacht (~/.cache/whisper) und
    beim nächsten Aufruf direkt verwendet.
    """
    ensure_ffmpeg()
    ensure("openai-whisper", "whisper")

    import whisper

    print(
        f"  Lade Whisper-Modell '{model_name}' "
        f"(nur beim ersten Mal – wird gecacht) …",
        file=sys.stderr,
    )
    model = whisper.load_model(model_name)

    print("  Transkribiere Audio …", file=sys.stderr)
    result = model.transcribe(str(p), verbose=False)

    # Ergebnis als strukturiertes Markdown aufbauen
    language = result.get("language", "unbekannt")
    segments = result.get("segments", [])

    lines = [f"*Erkannte Sprache: {language}*\n"]

    if segments:
        # Mit Zeitstempeln als Absätze formatieren
        for seg in segments:
            start = _format_ts(seg["start"])
            end   = _format_ts(seg["end"])
            text  = seg["text"].strip()
            lines.append(f"**[{start} → {end}]** {text}")
    else:
        lines.append(result.get("text", "").strip())

    return "\n\n".join(lines)


def _format_ts(seconds: float) -> str:
    """Formatiert Sekunden als MM:SS."""
    m, s = divmod(int(seconds), 60)
    return f"{m:02d}:{s:02d}"


# ─────────────────────────────────────────────────────────────────────────────
# Haupt-Konvertierungslogik
# ─────────────────────────────────────────────────────────────────────────────

def convert(
    input_path: Path,
    encoding: str = "utf-8",
    lang: str | None = None,
    model: str = "base",
    engine: str = "auto",
) -> str:
    ext = input_path.suffix.lower()

    # OCR-Sprache ermitteln (einmalig, lazy)
    ocr_lang = lang or best_ocr_lang()

    # — PDFs: Engine-Auswahl ------------------------------------------------
    if ext == ".pdf":
        use_marker = {
            "marker":  True,
            "classic": False,
            "auto":    marker_available(),
        }.get(engine)
        if use_marker is None:
            print(f"  Unbekannte Engine '{engine}' – nutze classic.", file=sys.stderr)
            use_marker = False

        if use_marker:
            try:
                return via_pdf_marker(input_path)
            except Exception as e:
                print(f"  marker-pdf: {e} → Fallback classic …", file=sys.stderr)
        return via_pdf(input_path, ocr_lang)

    # — markitdown (universell, versucht zuerst) ----------------------------
    try:
        ensure("markitdown[all]", "markitdown")
        from markitdown import MarkItDown
        result = MarkItDown().convert(str(input_path))
        content = (result.text_content or "").strip()
        if content:
            return content
    except Exception as e:
        print(f"  markitdown: {e} → Fallback …", file=sys.stderr)

    # — Format-spezifische Fallbacks ----------------------------------------
    dispatch = {
        ".csv":  via_csv,
        ".tsv":  via_csv,
        ".xlsx": via_xlsx,
        ".xls":  via_xlsx,
        ".html": via_html,
        ".htm":  via_html,
        ".docx": via_docx,
        ".doc":  via_docx,
        ".pptx": via_pptx,
        ".ppt":  via_pptx,
        ".txt":  lambda p: via_plain(p, encoding),
        ".md":   lambda p: via_plain(p, encoding),
        ".rst":  lambda p: via_plain(p, encoding),
        ".log":  lambda p: via_plain(p, encoding),
        # Bilddateien → direkt per OCR
        ".png":  lambda p: via_image_ocr(p, ocr_lang),
        ".jpg":  lambda p: via_image_ocr(p, ocr_lang),
        ".jpeg": lambda p: via_image_ocr(p, ocr_lang),
        ".tiff": lambda p: via_image_ocr(p, ocr_lang),
        ".tif":  lambda p: via_image_ocr(p, ocr_lang),
        ".bmp":  lambda p: via_image_ocr(p, ocr_lang),
        ".gif":  lambda p: via_image_ocr(p, ocr_lang),
        ".webp": lambda p: via_image_ocr(p, ocr_lang),
        # Audiodateien → Whisper-Transkription
        ".mp3":  lambda p: via_audio(p, model),
        ".wav":  lambda p: via_audio(p, model),
        ".m4a":  lambda p: via_audio(p, model),
        ".ogg":  lambda p: via_audio(p, model),
        ".flac": lambda p: via_audio(p, model),
        ".mp4":  lambda p: via_audio(p, model),   # Video → nur Audio
        ".mov":  lambda p: via_audio(p, model),
        ".mkv":  lambda p: via_audio(p, model),
    }

    if ext in dispatch:
        return dispatch[ext](input_path)

    # — Letzter Ausweg: als Plaintext lesen ----------------------------------
    print(f"  Kein Handler für {ext} – lese als Plaintext.", file=sys.stderr)
    return via_plain(input_path, encoding)


# ─────────────────────────────────────────────────────────────────────────────
# Nachbearbeitung
# ─────────────────────────────────────────────────────────────────────────────

def clean(text: str) -> str:
    """Normalisiert überschüssige Leerzeilen."""
    import re
    return re.sub(r"\n{4,}", "\n\n\n", text).strip()


def build_output(stem: str, source_name: str, body: str) -> str:
    """
    Baut den finalen Markdown-Text auf.
    Extrahiert den Dokumenttitel aus dem ersten # Heading im Body —
    so wird kein doppelter Titel erzeugt.
    """
    import re
    m = re.match(r"^# (.+)", body.lstrip())
    if m:
        title = m.group(1).strip()
        body  = re.sub(r"^\s*# .+\n*", "", body, count=1)
    else:
        title = stem.replace("_", " ").replace("-", " ").title()
    header = f"# {title}\n\n*Quelle: `{source_name}`*\n\n---\n\n"
    return header + body.lstrip()


# ─────────────────────────────────────────────────────────────────────────────
# Batch-Modus
# ─────────────────────────────────────────────────────────────────────────────

def batch_convert(input_dir: Path, output_dir: Path, **kwargs):
    """Konvertiert alle unterstützten Dateien in input_dir nach output_dir."""
    supported = {
        ".pdf", ".csv", ".tsv", ".xlsx", ".xls", ".html", ".htm",
        ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md", ".rst", ".log",
        ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif", ".webp",
    }
    files = sorted(p for p in input_dir.rglob("*") if p.suffix.lower() in supported)
    if not files:
        print(f"  Keine unterstützten Dateien in {input_dir}", file=sys.stderr)
        return 0

    output_dir.mkdir(parents=True, exist_ok=True)
    count = 0
    for p in files:
        rel = p.relative_to(input_dir)
        out = output_dir / rel.with_suffix(".md")
        out.parent.mkdir(parents=True, exist_ok=True)
        try:
            body   = convert(p, **kwargs)
            body   = clean(body)
            result = build_output(p.stem, p.name, body)
            out.write_text(result, encoding="utf-8")
            print(f"  ✓ {rel} → {out}", file=sys.stderr)
            count += 1
        except Exception as e:
            print(f"  ✗ {rel}: {e}", file=sys.stderr)
    return count


# ─────────────────────────────────────────────────────────────────────────────
# Einstiegspunkt
# ─────────────────────────────────────────────────────────────────────────────

def main():
    args     = sys.argv[1:]
    encoding = "utf-8"
    lang     = None
    model    = "base"
    engine   = "auto"

    def pop_flag(flag):
        if flag in args:
            idx = args.index(flag)
            val = args[idx + 1]
            args[:] = [a for i, a in enumerate(args) if i not in (idx, idx + 1)]
            return val
        return None

    encoding = pop_flag("--encoding") or encoding
    lang     = pop_flag("--lang")     or lang
    model    = pop_flag("--model")    or model
    engine   = pop_flag("--engine")   or engine

    batch = "--batch" in args
    if batch:
        args.remove("--batch")

    if not args:
        print(
            "Nutzung: convert_to_md.py <eingabe> [ausgabe.md]\n"
            "  [--engine marker|classic|auto]  (Standard: auto)\n"
            "  [--encoding <enc>] [--lang <tesseract-lang>] [--model <whisper>]\n"
            "  convert_to_md.py <ordner> --batch <ausgabe-ordner> [--engine ...]",
            file=sys.stderr,
        )
        sys.exit(1)

    input_path  = Path(args[0])
    output_path = Path(args[1]) if len(args) > 1 else None

    if not input_path.exists():
        print(f"Fehler: Datei nicht gefunden: {input_path}", file=sys.stderr)
        sys.exit(1)

    # — Batch-Modus ---------------------------------------------------------
    if batch:
        out_dir = output_path or Path.cwd() / "md_out"
        print(
            f"Batch: {input_path} → {out_dir} "
            f"(engine={engine}, {len(list(input_path.rglob('*')))} Dateien)",
            file=sys.stderr,
        )
        n = batch_convert(
            input_path, out_dir,
            encoding=encoding, lang=lang, model=model, engine=engine,
        )
        print(f"Fertig: {n} Dateien konvertiert.", file=sys.stderr)
        sys.exit(0 if n else 1)

    # — Einzeldatei ---------------------------------------------------------
    print(f"Konvertiere {input_path.name} … (engine={engine})", file=sys.stderr)
    body   = convert(input_path, encoding=encoding, lang=lang, model=model, engine=engine)
    body   = clean(body)
    result = build_output(input_path.stem, input_path.name, body)

    if output_path:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(result, encoding="utf-8")
        print(f"✓ Gespeichert → {output_path}", file=sys.stderr)
    else:
        print(result)


if __name__ == "__main__":
    main()
